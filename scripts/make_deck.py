#!/usr/bin/env python3
"""Generate the pre-final pitch deck as a real .pptx, with motion.

OWNER: Sorour.  Run from the repo root:

    .venv-main/bin/python scripts/make_deck.py

Writes docs/pitch/TheAgentOrg-prefinal.pptx. Committed and re-runnable rather than a
one-off, because somebody has to be able to fix a typo at 11pm on the 24th and
regenerate without reconstructing how the deck was built.

IT COVERS THE FOUR SECTIONS THE ORGANISER ASKED FOR
===================================================
From Hesham Khalil's invitation, and `_REQUIRED` below asserts each one is present as a
slide kicker so a rewrite cannot silently drop one:

    1. Overview of the idea (problem and proposed solution)  -> slides 4-7
    2. High-level architecture                               -> slides 8-10
       slide 8 is the DIAGRAM -- what runs where, as a picture
    3. Progress to date (what has been completed)            -> slides 11-13
    4. High-level future plan (roadmap to the final phase)   -> slide 14
    + a live demonstration                                   -> slide 15 hands over

Slide 2 is an agenda and slide 3 is the team -- both requested, and the agenda earns its
place by telling a judge up front that all four required sections are coming.

IT IS A PITCH, NOT A PROJECT REPORT
===================================
Three rules, and the first version of this deck broke all three:

  * NO SOURCE CODE ON A SLIDE. A judge cannot read five lines of Python off a Teams
    screen share in sixty seconds, and asking them to try spends their attention on
    parsing instead of on the argument. The idea is stated in English; the code is in
    the repository for anyone who wants it.
  * NO PER-PERSON SLIDES. Progress belongs to the team. A slide headed with one
    engineer's name invites "so what did the other four do", and makes a five-person
    team read as five people who worked separately.
  * NO "one function decides whether it ships -- and it has no model in it". It was on
    three slides. Cut, and `verify()` fails the build if it comes back.

TRANSITIONS AND ANIMATIONS ARE NOT A python-pptx FEATURE
========================================================
Measured before relying on it: `dir(slide)` exposes nothing matching "trans" or "anim".
python-pptx models shapes and text, not the timing tree. Both live in the slide's raw
XML, which it DOES expose, so `_transition` and `_animate` build that XML directly --
and `verify()` reads the saved archive back and counts them, because python-pptx would
never report them missing. It never knew about them.

MOTION IS ONE EFFECT, USED CONSISTENTLY. Content slides wipe in from the left; the four
statement slides fade. A different effect per slide is the single most reliable way to
make a deck look amateur. Entrance animations are click-advanced FADES -- never fly-ins
or spins, which read as a school project and pull the eye away from the words.

THE PALETTE
===========
DARK, and agreed in a browser before any of it reached this file. `docs/pitch/preview/`
renders the same geometry as HTML, so a colour decision took one reload instead of a
regenerate-and-open-PowerPoint cycle. Its CSS custom properties and the constants below
are the two places a colour lives; they must change together, or the thing on screen
stops being the thing that was signed off.

  VOID/SLATE  near-black surfaces, a hint of blue. Not flat grey, which reads as dead
  INK/DIM     off-white primary, muted secondary -- never pure white, which glares
  CYAN        every structural mark: kickers, rules, the pipeline, agent names
  ROSE/MINT   refused and shipped, used on the two slides that need them

Dark because the delivery medium is a Teams screen share, where the viewer's own
brightness is unknown, and because it reads as instrumentation rather than as a
corporate template.

PHOTOGRAPHS ARE PRE-CROPPED SQUARE
==================================
`docs/pitch/photos/square/` holds 640x640 JPEGs. CSS could crop on the fly with
`object-fit: cover`, but PowerPoint has no equivalent -- a non-square image in a square
frame is STRETCHED, and a stretched face is the one defect an audience notices instantly.
So the crop happens once, on disk, and both renderers show the same thing. Portraits are
biased toward the upper third, because a centre crop on a tall photograph cuts the top of
the head off.
"""

from __future__ import annotations

import itertools
import pathlib
import subprocess
import sys

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── verified numbers ──────────────────────────────────────────────────────────
# EVERY NUMBER HERE CARRIES THE COMMAND THAT PRODUCED IT, and all of them were
# RE-MEASURED for the final deck at `d6165c8` on 2026-08-28. That is not ceremony:
# every figure in the pre-final deck below was true when written and six of them were
# stale within a day, because twelve lanes committed concurrently. The test count moved
# 1102 -> 1853 and the file count 41 -> 73.
#
# Two of these are RANGES rather than point values, deliberately. CLAUDE.md records
# 116.88s -> 149.68s -> 102.83s for one unchanged test snapshot, purely load-dependent,
# so "measured" is a value PLUS its conditions and spread. A cost quoted as a point
# value would be the most repeated and least reproducible number on the deck.
TESTS_PASSING = 1853     # pytest -q | tail -1   (4 skip in a worktree, 3 on main)
TEST_FILES = 73          # ls tests/test_*.py | wc -l
WEB_TESTS = 166          # cd web && npm test   -> READ THE FILE COUNT: 10 files
WEB_TEST_FILES = 10      # same command. `Tests` counts what RAN; `Test Files` does not
TF_RESOURCES = 20        # grep -rhc '^resource ' infra/Terraform/modules/*/main.tf
AGENTORG_LINES = 20475   # find agentorg -name '*.py' | xargs wc -l | tail -1
WEB_LINES = 8769         # find web -name '*.ts' -o -name '*.tsx' | xargs wc -l
RUNTIME_VERSION = 37     # scripts/preflight.py check 2 — all five READY at the SAME one
CLEAN_MINUTES = 5        # measured, run 32585658981
POISONED_MINUTES = 3     # measured, run 32586453254
TRIGGER_SECONDS = 6      # issue created 16:45:09 -> run created 16:45:15

# scripts/measure_cost.py --runs 3 --require-model, at d6165c8. A RANGE: the spread is
# 30% across three consecutive runs of one unchanged ticket.
COST_LOW = "1.3"         # cents. $0.013036
COST_HIGH = "1.7"        # cents. $0.016931
COST_INFRA_SHARE = "0.1" # percent. $0.00001245 of $0.0131 is 0.095%

# The organiser's four required sections, as the kicker text that must appear.
_REQUIRED = ("THE PROBLEM", "THE SOLUTION", "ARCHITECTURE", "PROGRESS", "WHAT IS NEXT")

# ── palette ───────────────────────────────────────────────────────────────────
# PORTED VERBATIM FROM THE APPROVED PREVIEW, docs/pitch/preview/index.html. These hex
# values and that file's CSS custom properties are the two places a colour lives, and
# they must change together -- the preview is how the design was agreed, so a drift
# between them means the thing on screen is not the thing that was signed off.
#
# Dark, because a Teams screen share is the delivery medium: the viewer's brightness is
# unknown, and this reads as instrumentation rather than as a corporate template.
VOID = RGBColor(0x0A, 0x0C, 0x10)     # deepest black, for the statement slides
SLATE = RGBColor(0x14, 0x18, 0x1F)    # the standard slide surface -- near-black, a hint of blue
RAISED = RGBColor(0x1C, 0x22, 0x30)   # cards and panels sitting above the surface
LINE = RGBColor(0x2A, 0x32, 0x42)     # hairlines
INK = RGBColor(0xE8, 0xED, 0xF4)      # primary text: off-white, never pure white
DIM = RGBColor(0x8A, 0x94, 0xA6)      # secondary text
CYAN = RGBColor(0x4F, 0xD1, 0xC5)     # ALL structure: kickers, rules, the pipeline
ROSE = RGBColor(0xFF, 0x6B, 0x6B)     # the failure state
MINT = RGBColor(0x4A, 0xDE, 0x80)     # the success state
AMBER = RGBColor(0xF0, 0xA9, 0x3B)    # the human gates on the architecture diagram, and
                                      # nowhere else -- they are neither a success nor a
                                      # failure, they are a person deciding
SANS = "Helvetica Neue"
MONO = "Menlo"                        # agent names and identifiers: they are code, not labels

P_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Hoisted rather than written as `left=Inches(1.1)` defaults: ruff's B008 forbids a call
# in an argument default, and it is right to -- the value is evaluated once at import.
MARGIN = Inches(1.1)
BODY_WIDTH = Inches(11.1)
RULE_WIDTH = Inches(1.6)
STAT_WIDTH = Inches(3.4)      # the default width of a big-number block
ARROW_WIDTH = Inches(0.9)     # a connector arrow between diagram zones
ARROW_PAD = Inches(0.25)      # how far a connector's label may overhang its arrow


def _transition(slide, *, kind: str = "wipe", direction: str = "l") -> None:
    """Give this slide an entrance transition.

    Appended to <p:sld> as its LAST child, which the schema requires -- an element out
    of order makes PowerPoint declare the file corrupt and offer to repair it, which on
    a projector is indistinguishable from a broken deck.
    """
    element = etree.SubElement(slide._element, P_NS + "transition")
    element.set("spd", "med")
    child = etree.SubElement(element, P_NS + kind)
    if kind in ("wipe", "push", "pull", "cover"):
        child.set("dir", direction)


def _animate(slide, shape_ids: list[int]) -> None:
    """Reveal these shapes one click at a time, in the order given.

    ONE <p:timing> tree holding a sequence of click-triggered fades. Raw XML because
    there is no API for it, and built as a string rather than by element assembly: the
    nesting is five levels deep, and assembling that node by node would be far harder
    to read and to correct than the shape it produces.

    Each shape gets <p:set> to make it visible, then <p:animEffect filter="fade">. The
    <p:set> is what hides it beforehand -- without it the shape is on screen from the
    start and the fade animates something already visible.
    """
    if not shape_ids:
        return

    node_id = 10
    blocks = []
    for index, shape_id in enumerate(shape_ids):
        blocks.append(f"""
        <p:par><p:cTn id="{node_id}" fill="hold">
          <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
          <p:childTnLst><p:par><p:cTn id="{node_id + 1}" fill="hold">
            <p:stCondLst><p:cond delay="0"/></p:stCondLst>
            <p:childTnLst><p:par><p:cTn id="{node_id + 2}" presetID="10"
                presetClass="entr" presetSubtype="0" fill="hold"
                grpId="0" nodeType="{"clickEffect" if index == 0 else "afterEffect"}">
              <p:stCondLst><p:cond delay="0"/></p:stCondLst>
              <p:childTnLst>
                <p:set><p:cBhvr><p:cTn id="{node_id + 3}" dur="1" fill="hold"/>
                  <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                  <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                </p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>
                <p:animEffect transition="in" filter="fade">
                  <p:cBhvr><p:cTn id="{node_id + 4}" dur="450"/>
                    <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl></p:cBhvr>
                </p:animEffect>
              </p:childTnLst></p:cTn></p:par></p:childTnLst>
          </p:cTn></p:par></p:childTnLst>
        </p:cTn></p:par>""")
        node_id += 10

    timing = f"""<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
      <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst><p:seq concurrent="1" nextAc="seek">
          <p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>
            {"".join(blocks)}
          </p:childTnLst></p:cTn>
          <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
          <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
        </p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>"""
    slide._element.append(etree.fromstring(timing))


def _blank(prs, *, band=None):
    """A paper-coloured slide with no placeholders.

    Layout 6 is the blank one. The placeholder layouts fight explicit positioning -- a
    title placeholder re-centres itself and cannot be moved reliably across PowerPoint
    and Keynote, which is why every text box here is built by hand.

    `band` paints a colour block down the left edge: the deck's only recurring ornament,
    and what makes the statement slides feel like the same deck as the content ones.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    background.fill.solid()
    background.fill.fore_color.rgb = SLATE
    background.line.fill.background()
    background.shadow.inherit = False
    if band is not None:
        strip = slide.shapes.add_shape(1, 0, 0, Inches(0.34), SLIDE_H)
        strip.fill.solid()
        strip.fill.fore_color.rgb = band
        strip.line.fill.background()
        strip.shadow.inherit = False
    return slide


def _text(slide, text, *, left, top, width, height=None, size=20, color=INK,
          bold=False, font=SANS, align=PP_ALIGN.LEFT, spacing=1.2):
    """One text box. Returns the shape so its id can be animated."""
    box = slide.shapes.add_textbox(left, top, width, height or Inches(1))
    frame = box.text_frame
    frame.word_wrap = True
    for index, line in enumerate(str(text).split("\n")):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return box


def _rule(slide, *, top, left=MARGIN, width=RULE_WIDTH, color=CYAN):
    """A short accent rule under a heading. Cheap, and it makes a slide look designed."""
    bar = slide.shapes.add_shape(1, left, top, width, Emu(38100))  # ~0.04"
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def _heading(slide, text, *, kicker=None, size=40, top=None):
    """The standard slide head: optional kicker, title, accent rule."""
    if kicker:
        _text(slide, kicker.upper(), left=MARGIN, top=Inches(0.72), width=BODY_WIDTH,
              height=Inches(0.32), size=12, color=CYAN, bold=True, spacing=1.0)
    y = top or Inches(1.18)
    _text(slide, text, left=MARGIN, top=y, width=BODY_WIDTH, size=size, color=INK,
          bold=True, spacing=1.06)
    lines = str(text).count("\n") + 1
    _rule(slide, top=y + Inches(0.62 * lines + 0.30))


def _footer(slide, number):
    """Slide number, bottom right. Judges refer to slides by number when they ask.

    An explicit height, unlike most boxes: the 1in default would hang past the bottom
    edge and make the layout audit report a false positive on every slide. An audit that
    always warns is one nobody reads.
    """
    _text(slide, str(number), left=Inches(12.3), top=Inches(6.86), width=Inches(0.6),
          height=Inches(0.32), size=11, color=DIM, align=PP_ALIGN.RIGHT)


def _bullets(slide, items, *, top, size=20, gap=0.9, left=None, width=None, color=INK):
    """A stack of lines, each its own shape so each animates on its own click."""
    left = MARGIN if left is None else left
    width = BODY_WIDTH if width is None else width
    return [
        _text(slide, item, left=left, top=top + Inches(gap * index), width=width,
              size=size, color=color)
        for index, item in enumerate(items)
    ]


def _stat(slide, value, label, *, left, top, value_color=INK, width=STAT_WIDTH):
    """A big number with a caption under it. Two shapes, returned as a pair."""
    big = _text(slide, value, left=left, top=top, width=width, height=Inches(0.95),
                size=42, color=value_color, bold=True, spacing=1.0)
    small = _text(slide, label, left=left, top=top + Inches(0.88), width=width,
                  height=Inches(0.85), size=15, color=DIM, spacing=1.25)
    return [big, small]


# ── the slides ────────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = _blank(prs, band=CYAN)
    _text(slide, "THE AGENT ORG", left=MARGIN, top=Inches(2.25), width=BODY_WIDTH,
          size=64, color=INK, bold=True, spacing=1.0)
    _rule(slide, top=Inches(3.42), width=Inches(2.6))
    _text(slide, "AI agents that ship code the way an engineering team does —\n"
                 "with a safety check they cannot argue with.",
          left=MARGIN, top=Inches(3.82), width=Inches(10.4), size=24, color=DIM)
    _text(slide, "RosettaTeam   ·   DevOps Hackathon, final evaluation\n"
                 "28 August 2026",
          left=MARGIN, top=Inches(5.75), width=BODY_WIDTH, size=15, color=DIM)
    _transition(slide, kind="fade")
    _footer(slide, 1)


def slide_agenda(prs):
    """Slide 2. Tells a judge every section they asked for is coming.

    It earns its place for one reason: without it, a judge spends the first five minutes
    wondering whether we are going to cover the four topics the invitation named. With
    it, they can stop tracking that and listen.
    """
    slide = _blank(prs)
    _heading(slide, "Twenty minutes, five parts", kicker="agenda", size=40)
    heads = []
    y = Inches(2.6)
    for index, (title, detail) in enumerate([
        ("The team", "who built it"),
        ("The problem, and our solution",
         "why an AI pipeline needs a gate it cannot argue with"),
        ("Architecture", "how a ticket becomes a merged change"),
        ("Progress to date", "what runs today, what it costs, and what does not"),
        ("Roadmap, then a live demo", "two tickets: one ships, one is refused"),
    ], start=1):
        _text(slide, f"{index:02d}", left=MARGIN, top=y, width=Inches(0.7),
              height=Inches(0.5), size=17, color=CYAN, bold=True, font=MONO)
        heads.append(_text(slide, title, left=Inches(2.0), top=y, width=Inches(4.6),
                           height=Inches(0.5), size=20, color=INK, bold=True))
        _text(slide, detail, left=Inches(6.7), top=y, width=Inches(5.5),
              height=Inches(0.5), size=16, color=DIM)
        y += Inches(0.82)
    _transition(slide)
    _animate(slide, [s.shape_id for s in heads])
    _footer(slide, 2)


# The team, in speaking order. `photo` is a stem under docs/pitch/photos/square/, and
# `org` is deliberately blank-able: three of us have no employer to name here, and
# inventing one would be worse than the gap.
TEAM = [
    ("Mohamed Sorour", "Senior DevOps Engineer", "VEZEETA", "sorour"),
    ("Mariam Abdelkader", "Associate Solution Engineer", "RENOSYSTEMS", "mariam"),
    ("Habiba Megahed", "Junior DevOps Engineer", "DIGILIANS ALUM", "habiba"),
    ("Reem Shkeep", "Junior Testing Engineer", "DIGILIANS ALUM", "reem"),
    ("Aya Ebrahim", "Junior Testing Engineer", "DIGILIANS ALUM", "aya"),
]


def slide_team(prs):
    """Slide 3. Five members, a circular photograph above each.

    THE PHOTOGRAPHS ARE MASKED TO CIRCLES by setting `auto_shape_type = OVAL` on the
    picture, which writes `<a:prstGeom prst="ellipse">` into the drawing -- verified in
    the saved XML rather than assumed. A square photo in an oval frame is cropped by
    PowerPoint itself, so the sources are already square (see the module docstring) and
    nothing is stretched.

    A missing file degrades to initials rather than raising: the deck must still build
    for anyone who has not copied the photographs onto their machine.
    """
    slide = _blank(prs)
    _heading(slide, "RosettaTeam", kicker="meet the team", size=40)

    photos = pathlib.Path(__file__).resolve().parent.parent / "docs/pitch/photos/square"
    column = Inches(2.28)
    diameter = Inches(1.62)
    x = MARGIN
    shapes = []
    for name, role, org, stem in TEAM:
        centre = x + (column - diameter) / 2 - Inches(0.12)
        image = photos / f"{stem}.jpg"
        if image.exists():
            picture = slide.shapes.add_picture(str(image), centre, Inches(2.62),
                                               diameter, diameter)
            picture.auto_shape_type = MSO_SHAPE.OVAL
            picture.line.color.rgb = LINE
            picture.line.width = Pt(1.5)
            shapes.append(picture)
        else:
            ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, centre, Inches(2.62),
                                          diameter, diameter)
            ring.fill.solid()
            ring.fill.fore_color.rgb = RAISED
            ring.line.color.rgb = LINE
            ring.shadow.inherit = False
            para = ring.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = "".join(word[0] for word in name.split()[:2])
            run.font.size = Pt(26)
            run.font.bold = True
            run.font.color.rgb = CYAN
            run.font.name = MONO
            shapes.append(ring)

        _text(slide, name, left=x - Inches(0.16), top=Inches(4.48),
              width=column, height=Inches(0.62), size=16, color=INK, bold=True,
              align=PP_ALIGN.CENTER, spacing=1.15)
        _text(slide, role, left=x - Inches(0.16), top=Inches(5.08),
              width=column, height=Inches(0.66), size=13, color=DIM,
              align=PP_ALIGN.CENTER, spacing=1.2)
        if org.strip():
            _text(slide, org, left=x - Inches(0.16), top=Inches(5.74),
                  width=column, height=Inches(0.34), size=11, color=CYAN,
                  align=PP_ALIGN.CENTER, font=MONO)
        x += column

    _text(slide, "Three of us trained together at Digilians — this is the first thing "
                 "we have built as one team.",
          left=MARGIN, top=Inches(6.52), width=Inches(11.0), height=Inches(0.45),
          size=16, color=DIM)
    _transition(slide)
    _animate(slide, [s.shape_id for s in shapes])
    _footer(slide, 3)


def slide_problem(prs):
    slide = _blank(prs)
    _heading(slide, "AI writes code faster than anyone\ncan review it",
             kicker="the problem", size=36)
    body = _bullets(slide, [
        "Teams are already letting AI agents open pull requests.",
        "Review is the bottleneck — so the temptation is to trust and merge.",
        "We built that unchecked pipeline first, on purpose, to see what happens.",
    ], top=Inches(3.4), size=21, gap=0.72)
    stat = _stat(slide, "10 of 10",
                 "changes carrying hardcoded cloud credentials\nwere merged",
                 left=MARGIN, top=Inches(5.6), value_color=ROSE, width=Inches(4.0))
    _transition(slide)
    _animate(slide, [s.shape_id for s in body] + [s.shape_id for s in stat])
    _footer(slide, 4)


def slide_insight(prs):
    slide = _blank(prs, band=ROSE)
    _text(slide, "And nobody noticed.", left=MARGIN, top=Inches(2.0),
          width=BODY_WIDTH, size=52, color=INK, bold=True, spacing=1.0)
    _rule(slide, top=Inches(3.15), width=Inches(2.0), color=ROSE)
    lines = _bullets(slide, [
        "Every job passed. Every dashboard was green.",
        "The credential reached the main branch and the pipeline reported success.",
        "A check that never ran looks exactly like a check that passed.",
    ], top=Inches(3.6), size=23, gap=0.95, color=INK)
    _transition(slide, kind="fade")
    _animate(slide, [s.shape_id for s in lines])
    _footer(slide, 5)


def slide_solution(prs):
    slide = _blank(prs)
    _heading(slide, "So we gave the agents an organisation",
             kicker="the solution", size=38)
    _text(slide, "A ticket walks the same path it would in a real team.",
          left=MARGIN, top=Inches(2.5), width=BODY_WIDTH, height=Inches(0.5),
          size=20, color=DIM)

    # THE AGENT'S REAL NAME UNDER EACH STAGE. These are the runtime identifiers from
    # agentorg/github_ops.py:959, not labels invented for a slide -- a judge who later
    # reads the repository finds the same five words, and in mono because that is what
    # they are: identifiers.
    x = MARGIN
    for stage, agent in (("PLAN", "planner"), ("BUILD", "developer"),
                         ("REVIEW", "reviewer"), ("SCAN", "security"),
                         ("RELEASE", "sre")):
        card = slide.shapes.add_shape(1, x, Inches(3.15), Inches(1.95), Inches(0.98))
        card.fill.solid()
        card.fill.fore_color.rgb = RAISED
        card.line.color.rgb = CYAN
        card.line.width = Pt(1.25)
        card.shadow.inherit = False
        frame = card.text_frame
        frame.word_wrap = True
        head = frame.paragraphs[0]
        head.alignment = PP_ALIGN.CENTER
        run = head.add_run()
        run.text = stage
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = CYAN
        run.font.name = SANS
        sub = frame.add_paragraph()
        sub.alignment = PP_ALIGN.CENTER
        run = sub.add_run()
        run.text = agent
        run.font.size = Pt(11)
        run.font.color.rgb = DIM
        run.font.name = MONO
        x += Inches(2.28)

    # Aligned under gates 1, 2 and 3 -- after BUILD/REVIEW, after SCAN, and before
    # RELEASE. The previous version spaced the three labels evenly, which put them under
    # the wrong stages and quietly misdescribed where a human actually intervenes.
    for offset in (Inches(2.28), Inches(6.84), Inches(9.12)):
        _text(slide, "▲ human approval", left=MARGIN + offset - Inches(0.16),
              top=Inches(4.24), width=Inches(2.28), height=Inches(0.34),
              size=11, color=CYAN, align=PP_ALIGN.CENTER, font=MONO)

    rest = _bullets(slide, [
        "Five specialist agents: they plan, write, critique, scan and sign off.",
        "Three points where a named human must approve before anything moves.",
        "And a safety check that is ordinary arithmetic, not judgement.",
    ], top=Inches(4.95), size=20, gap=0.7)
    _transition(slide)
    _animate(slide, [s.shape_id for s in rest])
    _footer(slide, 6)


def slide_gate(prs):
    slide = _blank(prs, band=CYAN)
    _heading(slide, "The safety check cannot be persuaded",
             kicker="the solution", size=38)
    body = _bullets(slide, [
        "Real scanners read the change and report what they find.",
        "A fixed severity threshold decides: at or above it, the change stops.",
        "That decision is arithmetic. There is nothing in it to convince.",
    ], top=Inches(2.75), size=21, gap=0.8)
    proof = _text(slide,
        "We tried talking it out of a block. The reply insisted the change was safe and "
        "the scanners were wrong — and the change stayed blocked, because those words "
        "were never part of the decision.",
        left=MARGIN, top=Inches(5.4), width=Inches(10.7), height=Inches(1.1),
        size=19, color=CYAN)
    _transition(slide)
    _animate(slide, [s.shape_id for s in body] + [proof.shape_id])
    _footer(slide, 7)


def _zone(slide, label, *, left, top, width, height):
    """A dashed container marking one side of the boundary, with a label on its edge.

    Three of these carry the whole diagram: GitHub, AWS, and the agent runtimes. The
    boundary is the point -- "what runs where" is the most common question about this
    system, and a dashed box answers it faster than any sentence.

    The label sits ON the border with a slide-coloured backing box, which is how a
    hand-drawn diagram marks a region. There is no PowerPoint primitive for it; two
    shapes stacked is the whole trick.
    """
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    frame.fill.background()
    frame.line.color.rgb = LINE
    frame.line.width = Pt(1)
    frame.line.dash_style = 2      # MSO_LINE_DASH_STYLE.DASH
    frame.shadow.inherit = False
    frame.adjustments[0] = 0.03

    # A backing patch so the border appears to break behind the text, rather than
    # running through it.
    patch = slide.shapes.add_shape(1, left + Inches(0.18), top - Inches(0.11),
                                   Inches(0.06) * len(label) + Inches(0.28), Inches(0.22))
    patch.fill.solid()
    patch.fill.fore_color.rgb = SLATE
    patch.line.fill.background()
    patch.shadow.inherit = False

    _text(slide, label.upper(), left=left + Inches(0.22), top=top - Inches(0.14),
          width=Inches(3.2), height=Inches(0.26), size=9, color=DIM, bold=True)
    return frame


def _node(slide, title, subtitle, *, left, top, width, accent=None, small=False):
    """One box in the diagram. Returns it so it can be animated.

    `accent` recolours the border and title: CYAN for the pipeline itself, AMBER for a
    human gate. Everything else is a plain raised card, so the eye lands on the two
    things that matter.
    """
    height = Inches(0.46) if small else Inches(0.62)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RAISED
    box.line.color.rgb = accent or LINE
    box.line.width = Pt(1.25)
    box.shadow.inherit = False
    box.adjustments[0] = 0.08

    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.09)
    frame.margin_top = frame.margin_bottom = Inches(0.03)

    head = frame.paragraphs[0]
    head.alignment = PP_ALIGN.LEFT
    run = head.add_run()
    run.text = title
    run.font.size = Pt(11 if small else 12)
    run.font.bold = True
    run.font.color.rgb = accent or INK
    run.font.name = MONO if small else SANS

    if subtitle:
        sub = frame.add_paragraph()
        sub.alignment = PP_ALIGN.LEFT
        run = sub.add_run()
        run.text = subtitle
        run.font.size = Pt(9)
        run.font.color.rgb = DIM
        run.font.name = SANS
    return box


def _connector(slide, glyph, label, *, left, top, width=ARROW_WIDTH,
               pad=ARROW_PAD):
    """An arrow between two zones, with a mono label under it."""
    # The heights are explicit and tight because these two sit 0.26in apart: the 1in
    # default, or even 0.3in on the glyph, overlaps the label and the layout audit says so.
    _text(slide, glyph, left=left, top=top, width=width, height=Inches(0.3),
          size=17, color=CYAN, align=PP_ALIGN.CENTER)
    # `pad` widens the label box either side so a short arrow can still centre a long
    # word. It must be 0 where the arrow sits in a narrow gap between two zones --
    # otherwise the label reaches into the zone and collides with a box inside it.
    _text(slide, label, left=left - pad, top=top + Inches(0.42),
          width=width + pad * 2, height=Inches(0.22), size=9, color=DIM,
          align=PP_ALIGN.CENTER, font=MONO)


def slide_diagram(prs):
    """Slide 8. The architecture as a PICTURE, which the numbered walk cannot be.

    Asked for directly, and it was the right call: the list on the next slide describes
    the sequence but never shows the SHAPE -- two repositories, what AWS owns, five
    isolated agents, and three gates that belong to GitHub rather than to our code.

    Everything on it is real. The repository names, the five agent names and the service
    names are the ones in the code, so a judge who opens the project afterwards finds the
    same words.
    """
    slide = _blank(prs)
    _heading(slide, "What runs where", kicker="architecture", size=40)

    top = Inches(2.28)
    shapes = []

    # ── GitHub: both repositories, and the three gates that live there ──────────
    _zone(slide, "GitHub", left=MARGIN, top=top, width=Inches(3.5), height=Inches(4.24))
    shapes.append(_node(slide, "auth-service", "the target repo · a ticket is opened here",
                        left=MARGIN + Inches(0.22), top=top + Inches(0.34),
                        width=Inches(3.06)))
    shapes.append(_node(slide, "TheAgentOrg", "the pipeline · seven jobs",
                        left=MARGIN + Inches(0.22), top=top + Inches(1.08),
                        width=Inches(3.06), accent=CYAN))
    for index, label in enumerate(("▲  GATE 1 · plan approved",
                                   "▲  GATE 2 · change approved",
                                   "▲  GATE 3 · release approved")):
        shapes.append(_node(slide, label, "",
                            left=MARGIN + Inches(0.22),
                            top=top + Inches(1.98 + 0.66 * index),
                            width=Inches(3.06), accent=AMBER, small=True))

    # ── AWS ────────────────────────────────────────────────────────────────────
    aws_left = MARGIN + Inches(4.35)
    _zone(slide, "AWS", left=aws_left, top=top, width=Inches(3.1), height=Inches(2.72))
    for index, (name, note) in enumerate((
        ("Lambda", "verifies the signature first"),
        ("EventBridge", "routes it · dead-letter queue"),
        ("Bedrock", "the model the agents call"),
    )):
        shapes.append(_node(slide, name, note, left=aws_left + Inches(0.22),
                            top=top + Inches(0.34 + 0.76 * index), width=Inches(2.66)))

    # ── the five agents ────────────────────────────────────────────────────────
    agents_left = MARGIN + Inches(8.0)
    _zone(slide, "AgentCore · 5 runtimes", left=agents_left, top=top,
          width=Inches(2.95), height=Inches(4.24))
    for index, (agent, does) in enumerate((
        ("planner", "breaks the ticket down"),
        ("developer", "writes the change"),
        ("reviewer", "critiques it"),
        ("security", "three real scanners"),
        ("sre", "release readiness"),
    )):
        shapes.append(_node(slide, agent, does, left=agents_left + Inches(0.2),
                            top=top + Inches(0.34 + 0.7 * index), width=Inches(2.55),
                            accent=CYAN if agent == "security" else None, small=True))
    _text(slide, "one image · five tags\ndiffering only by role",
          left=agents_left + Inches(0.2), top=top + Inches(3.86), width=Inches(2.55),
          height=Inches(0.5), size=10, color=DIM)

    # ── the three hops between the zones ───────────────────────────────────────
    _connector(slide, "→", "webhook", left=MARGIN + Inches(3.58), top=top + Inches(0.44))
    _connector(slide, "←", "dispatch", left=MARGIN + Inches(3.58), top=top + Inches(1.18))
    # pad=0: this arrow lives in the 0.5in gap between AWS and the agents, so a widened
    # label would overlap the reviewer box at x=9.30. Measured.
    _connector(slide, "→", "invoke", left=MARGIN + Inches(7.42), top=top + Inches(1.18),
               width=Inches(0.66), pad=Emu(0))

    _text(slide, "No laptop anywhere in that path   ·   no long-lived cloud credentials"
                 "   ·   the three gates are GitHub's, not ours to bypass",
          left=MARGIN, top=Inches(6.78), width=Inches(11.2), height=Inches(0.4),
          size=14, color=CYAN)
    _transition(slide)
    # Animated by ZONE rather than by box: fifteen clicks would be unpresentable, and the
    # three zones are the actual argument -- GitHub, AWS, the agents.
    _animate(slide, [shapes[0].shape_id, shapes[5].shape_id, shapes[8].shape_id])
    _footer(slide, 8)


def slide_architecture(prs):
    slide = _blank(prs)
    _heading(slide, "Cloud-native, start to finish", kicker="architecture", size=38)

    steps = [
        ("A ticket is opened", "a GitHub issue — the only trigger"),
        ("Signed and verified", "authenticated before anything runs"),
        ("Routed to the pipeline", "an event bus, with a dead-letter queue"),
        ("Five agents, five runtimes", "each isolated, all from one build"),
        ("Three approval gates", "the run pauses until a human clicks"),
    ]
    heads = []
    y = Inches(2.5)
    for index, (title, detail) in enumerate(steps, start=1):
        _text(slide, f"{index}", left=MARGIN, top=y, width=Inches(0.5),
              height=Inches(0.5), size=19, color=CYAN, bold=True)
        heads.append(_text(slide, title, left=Inches(1.75), top=y, width=Inches(4.4),
                           height=Inches(0.5), size=19, color=INK, bold=True))
        _text(slide, detail, left=Inches(6.3), top=y, width=Inches(5.9),
              height=Inches(0.5), size=17, color=DIM)
        y += Inches(0.78)

    _text(slide, f"{TF_RESOURCES} infrastructure resources, all defined as code   ·   "
                 "no long-lived cloud credentials anywhere",
          left=MARGIN, top=Inches(6.55), width=Inches(11.2), height=Inches(0.4),
          size=15, color=CYAN)
    _transition(slide)
    _animate(slide, [s.shape_id for s in heads])
    _footer(slide, 9)


def slide_gates_detail(prs):
    slide = _blank(prs)
    _heading(slide, "The humans are not a formality",
             kicker="architecture", size=38)
    body = _bullets(slide, [
        "Each gate is a platform-level approval, not a step the pipeline can skip.",
        "Until someone clicks, the next stage does not exist — it is never queued.",
        "Refusing is recorded on the ticket: who refused, and when.",
    ], top=Inches(2.8), size=21, gap=0.84)
    note = _text(slide,
        "When the safety check stops a change, the stages after it are never created. "
        "There is no branch to take and no flag to flip — the refusal is structural.",
        left=MARGIN, top=Inches(5.5), width=Inches(10.7), height=Inches(1.0),
        size=19, color=CYAN)
    _transition(slide)
    _animate(slide, [s.shape_id for s in body] + [note.shape_id])
    _footer(slide, 10)


def slide_progress(prs):
    slide = _blank(prs)
    _heading(slide, "Working today, end to end", kicker="progress", size=38)
    left = _bullets(slide, [
        "The full pipeline runs in the cloud on every ticket.",
        "Three real scanners run inside the agents' own environment.",
        "Both outcomes verified this week against the live system.",
        "Every run leaves a readable record on the ticket it came from.",
    ], top=Inches(2.5), size=20, gap=0.74, width=Inches(6.2))
    stats = _stat(slide, f"{TESTS_PASSING}",
                  f"automated tests passing\nacross {TEST_FILES} files",
                  left=Inches(8.0), top=Inches(2.5))
    stats += _stat(slide, f"~{CLEAN_MINUTES} min",
                   "from ticket opened\nto change merged",
                   left=Inches(8.0), top=Inches(4.5))
    _transition(slide)
    _animate(slide, [s.shape_id for s in left] + [s.shape_id for s in stats])
    _footer(slide, 11)


def slide_platform(prs):
    """Slide 12 — what the final phase added, and it is the slide that did not exist.

    THE PRE-FINAL DECK PREDATES ALL OF IT. Twelve lanes landed a durable queue, tenancy,
    deterministic scoring, cost instrumentation, a self-hosted path, integration
    adapters, a control-plane API, generated tests, retrieval and a web product. A deck
    that showed the same six progress slides would be describing last month's system.

    ONE LINE PER CAPABILITY AND NOT ONE MORE. Ten items at 17pt is already the densest
    slide in this deck; a second clause on each would make it a document. The bullets
    say what a thing DOES rather than what it is called, because a judge does not know
    what "Lane H" is and should not have to.
    """
    slide = _blank(prs)
    _heading(slide, "From a pipeline to a platform", kicker="progress", size=38)
    _text(slide, "The pre-final demo was one pipeline on one repository.",
          left=MARGIN, top=Inches(2.35), width=Inches(11.0), height=Inches(0.4),
          size=18, color=DIM)
    left = _bullets(slide, [
        "Runs queue, and a pause survives a restart.",
        "Many customers on one deployment, with a suite that tries to breach it.",
        "One severity table for three scanners, plus an audit row per finding.",
        "Every run reports what it cost, at the published rate.",
        "The same pipeline runs on your own hardware.",
    ], top=Inches(3.05), size=17, gap=0.72, width=Inches(5.5))
    right = _bullets(slide, [
        "GitHub is one adapter, not the foundation.",
        "A machine API to submit, watch and cancel — it cannot open a gate.",
        "Tests written from the ticket, never from the change.",
        "Agents read past decisions; that text cannot reach the verdict.",
        "A web app: watch a run live, approve or refuse.",
    ], top=Inches(3.05), size=17, gap=0.72, left=Inches(7.0), width=Inches(5.3))
    note = _text(slide,
        f"{AGENTORG_LINES:,} lines of pipeline and {WEB_LINES:,} of application, "
        f"under {TESTS_PASSING} tests and {WEB_TESTS} more for the web.",
        left=MARGIN, top=Inches(6.78), width=Inches(11.0), height=Inches(0.42),
        size=16, color=CYAN)
    _transition(slide)
    _animate(slide, [s.shape_id for s in left] + [s.shape_id for s in right]
             + [note.shape_id])
    _footer(slide, 12)


def slide_evidence(prs):
    slide = _blank(prs)
    _heading(slide, "The same request, twice", kicker="progress", size=38)
    _text(slide, "One ticket is clean. One carries a hardcoded credential. "
                 "Nothing else differs.",
          left=MARGIN, top=Inches(2.45), width=Inches(11.0), height=Inches(0.5),
          size=20, color=DIM)

    good = _text(slide, "SHIPPED", left=MARGIN, top=Inches(3.3), width=Inches(5.0),
                 height=Inches(0.58), size=26, color=MINT, bold=True)
    good_body = _text(slide,
        "Planned, written, reviewed, scanned,\n"
        "approved three times, merged.\n"
        "The ticket closed itself.",
        left=MARGIN, top=Inches(4.0), width=Inches(5.0), height=Inches(1.7),
        size=18, color=INK)

    bad = _text(slide, "REFUSED", left=Inches(7.2), top=Inches(3.3), width=Inches(4.9),
                height=Inches(0.58), size=26, color=ROSE, bold=True)
    bad_body = _text(slide,
        "Two credentials found in the change.\n"
        "Stopped before anyone could approve it,\n"
        "and the ticket says why.",
        left=Inches(7.2), top=Inches(4.0), width=Inches(4.9), height=Inches(1.7),
        size=18, color=INK)

    close = _text(slide,
        f"About {POISONED_MINUTES} minutes to refuse. Nothing merged, and the reason is "
        "on the ticket in plain English.",
        left=MARGIN, top=Inches(6.35), width=Inches(11.0), height=Inches(0.5),
        size=18, color=DIM)
    _transition(slide)
    _animate(slide, [good.shape_id, good_body.shape_id, bad.shape_id, bad_body.shape_id,
                     close.shape_id])
    _footer(slide, 13)


def slide_why_it_matters(prs):
    slide = _blank(prs, band=MINT)
    _heading(slide, "What this is worth", kicker="progress", size=40)
    body = _bullets(slide, [
        "Review capacity stops being the limit on how fast a team can ship.",
        "The checks that matter run on every change, not when someone remembers.",
        "A refusal is auditable: what was found, who approved, and when.",
        "It bolts onto the tools a team already uses — issues, branches, approvals.",
    ], top=Inches(2.85), size=21, gap=0.88)
    _transition(slide, kind="fade")
    _animate(slide, [s.shape_id for s in body])
    _footer(slide, 14)


def slide_cost(prs):
    """Slide 15 — the price of running every check on every change.

    THE NUMBER IS SMALL AND THAT IS THE ARGUMENT. A judge's unspoken objection to
    running five agents on every ticket is that it must be expensive. It is one and a
    half cents, measured, and the infrastructure around it is three orders of magnitude
    below the model.

    THE THIRD COLUMN IS DELIBERATELY EMPTY, and saying so out loud is worth more than
    filling it. A developer's hourly rate is not in this repository; a plausible figure
    would become the most quoted number in the room and the least defensible. This slide
    states the machine cost and hands the wage back to the audience, which is also the
    only honest way to present it.
    """
    slide = _blank(prs)
    _heading(slide, "What it costs to check everything", kicker="progress", size=38)
    _text(slide, "Measured over three runs of one ticket, priced from the published "
                 "rate card on the day it was read.",
          left=MARGIN, top=Inches(2.35), width=Inches(11.0), height=Inches(0.44),
          size=18, color=DIM)

    stats = _stat(slide, f"{COST_LOW}–{COST_HIGH}¢",
                  "per change, on the cloud —\nevery check, every agent",
                  left=MARGIN, top=Inches(3.2), value_color=MINT, width=Inches(3.5))
    stats += _stat(slide, f"{COST_INFRA_SHARE}%",
                   "of that is infrastructure.\nThe rest is the model",
                   left=Inches(4.9), top=Inches(3.2), width=Inches(3.5))
    stats += _stat(slide, "zero",
                   "on your own hardware —\nslower, same verdict",
                   left=Inches(8.6), top=Inches(3.2), width=Inches(3.5))

    note = _text(slide,
        "We do not put a figure on the third comparison — a developer's hour. That "
        "number is not ours to invent, and it would be the one everybody remembered.",
        left=MARGIN, top=Inches(5.5), width=Inches(11.0), height=Inches(0.9),
        size=18, color=CYAN)
    _transition(slide)
    _animate(slide, [s.shape_id for s in stats] + [note.shape_id])
    _footer(slide, 15)


def slide_limits(prs):
    """Slide 16 — the limitations, volunteered.

    THE PRE-FINAL EVALUATION'S STRONGEST MOMENT WAS VOLUNTEERING A LIMITATION BEFORE
    BEING ASKED, and this slide is that, on purpose, immediately before the roadmap so
    the roadmap reads as a response rather than as a wish list.

    EACH ONE IS COSTED IN THE REPOSITORY, and the closing line says where. A limitation
    merely admitted is weaker than one costed: the costing is what shows it is understood
    rather than noticed.

    THE HONEST SUMMARY IS THE LAST BULLET and it is the one to say aloud — the pattern is
    "unexecuted", not "broken", and that distinction is what the whole project is about.
    """
    slide = _blank(prs, band=AMBER)
    _heading(slide, "What does not work yet", kicker="progress", size=38)
    body = _bullets(slide, [
        "Sign-in has never completed — it needs a database nobody has started here.",
        "A run does not yet record its own cost, though the meter works.",
        "The queue that would replace our CI is tested and not yet switched on.",
        "The browser tests are written and no browser has run them.",
        "Any of our three approvals can be bypassed by a repository administrator.",
    ], top=Inches(2.7), size=19, gap=0.72)
    note = _text(slide,
        "Sixteen limitations are written down, each with what removing it would take. "
        "Every one is something unexecuted rather than something broken — and telling "
        "those two apart is the whole point of this project.",
        left=MARGIN, top=Inches(6.15), width=Inches(11.0), height=Inches(0.9),
        size=17, color=INK)
    _transition(slide, kind="fade")
    _animate(slide, [s.shape_id for s in body] + [note.shape_id])
    _footer(slide, 16)


def slide_roadmap(prs):
    slide = _blank(prs)
    _heading(slide, "Roadmap to the final phase", kicker="what is next", size=38)
    items = _bullets(slide, [
        "TURN ON WHAT IS BUILT — the queue, the database, and cost on every run.",
        "CACHE THE CONTEXT — the agents re-read the same code five times and pay for it.",
        "WIDEN THE CHECKS — dependency and licence scanning, per-project thresholds.",
        "PROVE IT AT VOLUME — ten times the sample behind every number in this deck.",
    ], top=Inches(2.6), size=19, gap=1.08)
    note = _text(slide,
        "Every item is a gap we have already written down and costed. We would rather "
        "show you a known limitation than discover one on stage.",
        left=MARGIN, top=Inches(6.5), width=Inches(11.0), height=Inches(0.5),
        size=17, color=CYAN)
    _transition(slide)
    _animate(slide, [s.shape_id for s in items] + [note.shape_id])
    _footer(slide, 17)


def slide_demo(prs):
    slide = _blank(prs, band=CYAN)
    _text(slide, "Let us show you.", left=MARGIN, top=Inches(2.35), width=BODY_WIDTH,
          size=52, color=INK, bold=True, spacing=1.0)
    _rule(slide, top=Inches(3.5), width=Inches(2.0))
    _text(slide, "Two tickets. The same feature request.\n"
                 "One ships itself. One is refused.",
          left=MARGIN, top=Inches(3.9), width=Inches(10.4), size=26, color=DIM)
    _transition(slide, kind="fade")
    _footer(slide, 18)


def slide_close(prs):
    slide = _blank(prs, band=CYAN)
    _heading(slide, "Thank you", size=44)
    close = _text(slide,
        "Agents did the work.\n"
        "Humans stayed in control of what shipped.\n"
        "And the change that should not have shipped, did not.",
        left=MARGIN, top=Inches(2.6), width=Inches(11.0), size=27, color=INK,
        spacing=1.55)
    team = _text(slide, "RosettaTeam   ·   questions welcome",
                 left=MARGIN, top=Inches(5.85), width=Inches(11.0), height=Inches(0.5),
                 size=17, color=DIM)
    _transition(slide, kind="fade")
    _animate(slide, [close.shape_id, team.shape_id])
    _footer(slide, 19)


SLIDES = [
    slide_title, slide_agenda, slide_team, slide_problem, slide_insight, slide_solution, slide_gate,
    slide_diagram, slide_architecture, slide_gates_detail, slide_progress, slide_platform,
    slide_evidence, slide_why_it_matters, slide_cost, slide_limits, slide_roadmap,
    slide_demo, slide_close,
]


def build(out: pathlib.Path) -> pathlib.Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    for make in SLIDES:
        make(prs)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return out


def _layout_problems(prs) -> list[str]:
    """Boxes whose WRAPPED text collides with the next box or runs off the slide.

    THE QUESTION IS WRAPPED HEIGHT, NOT LINE WIDTH. `word_wrap` is on, so a long line
    does not overflow sideways -- it wraps, and the box grows DOWNWARD into whatever
    sits below it. A width-only check reports nothing while three slides overlap, which
    is how the first version of this deck passed its own audit with six collisions.

    Estimated rather than measured, because the real answer needs a font renderer: line
    count from an average glyph advance, leading 1.25, plus the frame's insets.
    Approximate on purpose -- it is a smoke alarm, and a 0.1in collision is what it is
    for.

    Boxes are compared only when they OVERLAP HORIZONTALLY. Two columns side by side
    share a vertical band by design, and flagging those made this cry wolf on every
    stat slide.
    """
    problems = []
    for number, slide in enumerate(prs.slides, 1):
        boxes = []
        for shape in slide.shapes:
            if not shape.has_text_frame or not shape.text_frame.text.strip():
                continue
            lines, biggest = 0, 0
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs)
                if not text.strip():
                    lines += 1
                    continue
                pt = max((r.font.size.pt if r.font.size else 18) for r in para.runs)
                biggest = max(biggest, pt)
                per_line = max(1, int(shape.width / Inches(1) / (pt * 0.50 / 72)))
                lines += max(1, -(-len(text) // per_line))
            needed = Emu(int(lines * biggest * 1.25 * 12700)) + Inches(0.1)
            boxes.append((shape.top, needed, lines, biggest,
                          shape.text_frame.text[:42], shape.left,
                          shape.left + shape.width))
            if shape.top + needed > SLIDE_H:
                problems.append(
                    f"slide {number}: {shape.text_frame.text[:32]!r} runs "
                    f"{(shape.top + needed - SLIDE_H) / Inches(1):.2f}in past the bottom"
                )
        boxes.sort()
        for upper, lower in itertools.pairwise(boxes):
            top, needed, lines, pt, label, left, right = upper
            next_top, _n, _l, _p, _lb, next_left, next_right = lower
            if right <= next_left or next_right <= left:
                continue
            if top + needed > next_top:
                problems.append(
                    f"slide {number}: {label!r} ({lines} lines @{pt:.0f}pt) overlaps the "
                    f"next box by {(top + needed - next_top) / Inches(1):.2f}in"
                )
    return problems


def verify(path: pathlib.Path) -> int:
    """Assert the motion, the content rules and the layout reached the saved file.

    THE MOTION CHECK EARNS ITS PLACE. A deck that silently lost its transitions or its
    timing tree is byte-different but visually identical until it is presented, and
    python-pptx will not complain -- it never knew about them. So the archive is read
    back and the elements are counted.
    """
    import zipfile

    problems = []
    with zipfile.ZipFile(path) as archive:
        slides = sorted(n for n in archive.namelist()
                        if n.startswith("ppt/slides/slide") and n.endswith(".xml"))
        if len(slides) != len(SLIDES):
            problems.append(f"{len(slides)} slides in the file, expected {len(SLIDES)}")
        without_transition = []
        with_timing = 0
        for name in slides:
            xml = archive.read(name).decode("utf-8")
            if "<p:transition" not in xml:
                without_transition.append(name.rsplit("/", 1)[-1])
            if "animEffect" in xml:
                with_timing += 1
        if without_transition:
            problems.append(f"no transition on: {', '.join(without_transition)}")
        if with_timing < 9:
            problems.append(f"only {with_timing} slides carry entrance animations")

    text = " ".join(
        shape.text_frame.text
        for slide in Presentation(path).slides
        for shape in slide.shapes
        if shape.has_text_frame
    )

    # BANNED, and enforced rather than remembered: each was cut deliberately and each is
    # easy to reintroduce while editing prose.
    for phrase in ("no model in it", "def ", "return (", "compute_security_verdict"):
        if phrase in text:
            problems.append(f"banned phrase is back on a slide: {phrase!r}")

    # THE ORGANISER'S FOUR SECTIONS. A deck that reads beautifully and omits one of them
    # fails the brief, and the omission is invisible while writing.
    missing = [section for section in _REQUIRED if section not in text.upper()]
    if missing:
        problems.append(f"required section absent: {', '.join(missing)}")

    layout = _layout_problems(Presentation(path))
    problems.extend(layout)

    print(f"{path}  ({path.stat().st_size // 1024} KB)")
    print(f"  slides:      {len(SLIDES)}")
    print(f"  animated:    {with_timing}")
    print(f"  layout:      {'clean' if not layout else f'{len(layout)} collisions'}")
    print(f"  sections:    {'all four covered' if not missing else 'MISSING ' + str(missing)}")
    print("  transitions: all" if not without_transition else f"  MISSING: {without_transition}")
    if problems:
        for problem in problems:
            print(f"  FAIL: {problem}", file=sys.stderr)
        return 1
    print("  OK — motion, content rules and layout verified in the saved file")
    return 0


def main() -> int:
    root = pathlib.Path(__file__).resolve().parent.parent
    out = build(root / "docs" / "pitch" / "TheAgentOrg-final.pptx")
    code = verify(out)
    # `file(1)` is the independent witness: it reads the archive's own magic rather than
    # trusting the library that wrote it.
    kind = subprocess.run(["file", "-b", str(out)], capture_output=True, text=True,
                          check=False).stdout.strip()
    print(f"  file(1):     {kind}")
    return code


if __name__ == "__main__":
    sys.exit(main())
